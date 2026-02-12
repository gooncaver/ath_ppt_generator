"""
Slide Exporter - Export PowerPoint slides as images
Enables visual review by LLM during generation
"""

import os
import tempfile
from pathlib import Path
from typing import Optional, Tuple
from pptx import Presentation
from PIL import Image
import io


class SlideExporter:
    """Export individual slides as images for visual review"""
    
    def __init__(self, presentation_path: Optional[str] = None):
        """
        Initialize slide exporter
        
        Args:
            presentation_path: Optional path to existing presentation
        """
        self.presentation_path = presentation_path
        self.prs = Presentation(presentation_path) if presentation_path else None
        
    def export_slide(
        self,
        slide_index: int,
        output_path: Optional[str] = None,
        format: str = 'PNG',
        size: Tuple[int, int] = (1280, 720)
    ) -> str:
        """
        Export a single slide as an image
        
        Args:
            slide_index: Zero-based index of slide to export
            output_path: Optional output path (temp file if None)
            format: Image format (PNG, JPG)
            size: Output size (width, height)
            
        Returns:
            Path to exported image file
            
        Note:
            This uses python-pptx to save and external conversion.
            For production, consider PowerPoint COM automation on Windows
            or LibreOffice headless for better quality.
        """
        if not self.prs:
            raise ValueError("No presentation loaded")
        
        if slide_index >= len(self.prs.slides):
            raise IndexError(f"Slide index {slide_index} out of range (max {len(self.prs.slides)-1})")
        
        # Generate output path if not provided
        if output_path is None:
            temp_dir = Path(tempfile.gettempdir()) / "ath_ppt_generator"
            temp_dir.mkdir(exist_ok=True)
            output_path = str(temp_dir / f"slide_{slide_index}.{format.lower()}")
        
        # Use Windows COM automation if available (best quality)
        if os.name == 'nt':
            try:
                return self._export_with_com(slide_index, output_path, format, size)
            except ImportError:
                pass  # Fall back to python-pptx method
        
        # Fallback: Create single-slide presentation and convert
        return self._export_with_pptx(slide_index, output_path, format, size)
    
    def _export_with_com(
        self,
        slide_index: int,
        output_path: str,
        format: str,
        size: Tuple[int, int]
    ) -> str:
        """
        Export using Windows COM automation (PowerPoint application)
        Best quality but Windows-only
        """
        import win32com.client
        import pythoncom
        
        # Initialize COM
        pythoncom.CoInitialize()
        
        try:
            # Get absolute paths
            prs_path = str(Path(self.presentation_path).absolute())
            out_path = str(Path(output_path).absolute())
            
            # Start PowerPoint
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = 1  # Must be visible for export to work
            
            # Open presentation
            presentation = powerpoint.Presentations.Open(
                prs_path,
                ReadOnly=True,
                Untitled=True,
                WithWindow=False
            )
            
            # Export specific slide
            slide = presentation.Slides(slide_index + 1)  # COM uses 1-based indexing
            
            # Export as image
            slide.Export(out_path, format.upper())
            
            # Cleanup
            presentation.Close()
            powerpoint.Quit()
            
            return out_path
            
        finally:
            pythoncom.CoUninitialize()
    
    def _export_with_pptx(
        self,
        slide_index: int,
        output_path: str,
        format: str,
        size: Tuple[int, int]
    ) -> str:
        """
        Export using python-pptx (fallback method)
        Lower quality but cross-platform
        
        Note: This creates a temporary single-slide presentation
        and requires external conversion tool (like LibreOffice)
        """
        # Create new presentation with just this slide
        from pptx.util import Inches
        
        temp_prs = Presentation()
        temp_prs.slide_width = self.prs.slide_width
        temp_prs.slide_height = self.prs.slide_height
        
        # Copy slide layout
        source_slide = self.prs.slides[slide_index]
        slide_layout = source_slide.slide_layout
        
        # This is a simplified copy - may not preserve all formatting
        # For production use, implement full slide cloning
        
        # Save temp presentation
        temp_path = output_path.replace(f'.{format.lower()}', '_temp.pptx')
        temp_prs.save(temp_path)
        
        # Convert using external tool (requires LibreOffice or similar)
        # For now, return placeholder message
        # TODO: Implement conversion with LibreOffice headless or similar
        
        print(f"WARNING: Cross-platform export not fully implemented.")
        print(f"Saved temporary presentation: {temp_path}")
        print(f"Please convert manually to: {output_path}")
        
        return temp_path
    
    def export_all_slides(
        self,
        output_dir: str,
        format: str = 'PNG',
        size: Tuple[int, int] = (1280, 720)
    ) -> list:
        """
        Export all slides in presentation
        
        Args:
            output_dir: Directory to save slides
            format: Image format
            size: Output size
            
        Returns:
            List of exported file paths
        """
        if not self.prs:
            raise ValueError("No presentation loaded")
        
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        
        exported_files = []
        
        for i in range(len(self.prs.slides)):
            slide_path = output_path / f"Slide{i+1}.{format.upper()}"
            try:
                exported = self.export_slide(i, str(slide_path), format, size)
                exported_files.append(exported)
                print(f"  Exported slide {i+1}/{len(self.prs.slides)}: {slide_path.name}")
            except Exception as e:
                print(f"  Error exporting slide {i+1}: {e}")
        
        return exported_files
    
    @staticmethod
    def export_slide_from_presentation(
        prs: Presentation,
        slide_index: int,
        output_path: Optional[str] = None,
        format: str = 'PNG'
    ) -> str:
        """
        Export slide directly from a Presentation object
        
        Args:
            prs: python-pptx Presentation object
            slide_index: Zero-based slide index
            output_path: Optional output path
            format: Image format
            
        Returns:
            Path to exported image
        """
        # Save to temporary file first
        temp_pptx = Path(tempfile.gettempdir()) / f"temp_slide_{slide_index}.pptx"
        prs.save(str(temp_pptx))
        
        # Create exporter and export
        exporter = SlideExporter(str(temp_pptx))
        result = exporter.export_slide(slide_index, output_path, format)
        
        # Cleanup temp file
        try:
            temp_pptx.unlink()
        except:
            pass
        
        return result


def main():
    """CLI entry point for testing"""
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python slide_exporter.py <presentation.pptx> [output_dir]")
        sys.exit(1)
    
    prs_path = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "output/slide_exports"
    
    print(f"Exporting slides from: {prs_path}")
    print(f"Output directory: {output_dir}")
    
    exporter = SlideExporter(prs_path)
    exported = exporter.export_all_slides(output_dir)
    
    print(f"\n✓ Exported {len(exported)} slides")
    print(f"  Directory: {output_dir}")


if __name__ == "__main__":
    main()
