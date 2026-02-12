"""
Smart Slide Generator - Phase 2
Uses LLM intelligence for layout selection and content organization
"""

import sys
from pathlib import Path

# Add parent directory to path for imports when running as script
if __name__ == "__main__":
    sys.path.insert(0, str(Path(__file__).parent.parent))

from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
from typing import List, Dict, Optional, Any

from src.llm.client import LLMClient
from src.llm.planner import ContentPlanner
from src.template_inspector import TemplateInspector


class SmartGenerator:
    """AI-powered slide generator with intelligent layout selection"""
    
    def __init__(self, template_path: str, api_key: Optional[str] = None):
        """
        Initialize smart generator
        
        Args:
            template_path: Path to PowerPoint template file
            api_key: Optional OpenAI API key
        """
        self.template_path = Path(template_path)
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")
        
        # Initialize LLM components
        self.llm_client = LLMClient(api_key=api_key)
        self.planner = ContentPlanner(self.llm_client)
        
        # Load template
        self.prs = Presentation(str(self.template_path))
        
        # Get available layouts
        inspector = TemplateInspector(str(self.template_path))
        self.layouts = inspector.get_slide_layouts()
        all_layout_names = [layout['name'] for layout in self.layouts]
        
        # Filter to consistent background color (prefer dark backgrounds)
        # Exclude layouts with 'light', 'white', or 'Light' in the name
        dark_layouts = [name for name in all_layout_names 
                       if not any(keyword in name.lower() for keyword in ['light', 'white'])]
        
        # Use dark layouts if we have enough, otherwise use all
        self.layout_names = dark_layouts if len(dark_layouts) >= 10 else all_layout_names
        
        print(f"✓ Loaded template with {len(all_layout_names)} layouts (using {len(self.layout_names)} consistent layouts)")
    
    def get_layout_by_name(self, name: str):
        """
        Find layout by name in template
        
        Args:
            name: Layout name to search for
            
        Returns:
            Layout object or None
        """
        for layout in self.prs.slide_layouts:
            if layout.name == name:
                return layout
        return None
    
    def create_slide(
        self,
        layout_name: str,
        title: str,
        content: List[str],
        notes: str = ""
    ):
        """
        Create a slide with specified layout and content
        
        Args:
            layout_name: Name of layout to use
            title: Slide title
            content: List of content items/bullets
            notes: Presenter notes
        """
        # Get layout
        layout = self.get_layout_by_name(layout_name)
        if layout is None:
            # Fallback to a default content layout
            fallback_name = "10_Title and Content"
            layout = self.get_layout_by_name(fallback_name)
            if layout is None:
                layout = self.prs.slide_layouts[3] if len(self.prs.slide_layouts) > 3 else self.prs.slide_layouts[0]
            print(f"  Warning: Layout '{layout_name}' not found, using '{layout.name}'")
        
        # Create slide
        slide = self.prs.slides.add_slide(layout)
        
        # Populate title
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # Populate content
        if content:
            # Find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                # Look for body/content placeholder (usually idx 1)
                if shape.placeholder_format.idx == 1:
                    content_placeholder = shape
                    break
            
            if content_placeholder and hasattr(content_placeholder, 'text_frame'):
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                # Add content items as paragraphs
                for idx, item in enumerate(content):
                    if idx == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = item
                    p.level = 0
        
        # Add notes
        if notes and hasattr(slide, 'notes_slide'):
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = notes
    
    def generate_from_text(
        self,
        text_content: str,
        output_path: Optional[str] = None,
        target_slides: Optional[int] = None
    ) -> str:
        """
        Generate presentation from text content using AI
        
        Args:
            text_content: Raw text/markdown content
            output_path: Path to save output file (auto-generated if None)
            target_slides: Optional target number of slides
            
        Returns:
            Path to generated presentation
        """
        print("\n" + "="*70)
        print("SMART SLIDE GENERATION (Phase 2)")
        print("="*70 + "\n")
        
        # Clear existing template slides
        existing_count = len(self.prs.slides)
        if existing_count > 0:
            print(f"Removing {existing_count} existing slides from template...")
            for i in range(existing_count - 1, -1, -1):
                rId = self.prs.slides._sldIdLst[i].rId
                self.prs.part.drop_rel(rId)
                del self.prs.slides._sldIdLst[i]
        
        # Create slide plan with AI
        print(f"\nAvailable layouts: {len(self.layout_names)}")
        slide_plan = self.planner.create_slide_plan(
            content=text_content,
            available_layouts=self.layout_names,
            target_slides=target_slides
        )
        
        if not slide_plan:
            print("Error: No slide plan generated")
            return None
        
        print(f"\n{'─'*70}")
        print(f"Generating {len(slide_plan)} slides...")
        print(f"{'─'*70}\n")
        
        # Create slides from plan
        for slide_spec in slide_plan:
            slide_num = slide_spec.get('slide_number', 0)
            title = slide_spec.get('title', '')
            layout_name = slide_spec.get('layout_name', '')
            content = slide_spec.get('content', [])
            notes = slide_spec.get('notes', '')
            
            print(f"  [{slide_num}/{len(slide_plan)}] {title}")
            print(f"      Layout: {layout_name}")
            if content:
                print(f"      Bullets: {len(content)}")
            
            self.create_slide(
                layout_name=layout_name,
                title=title,
                content=content,
                notes=notes
            )
        
        # Save presentation
        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = f"output/smart_presentation_{timestamp}.pptx"
        
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        self.prs.save(str(output_path))
        
        # Print summary
        print(f"\n{'='*70}")
        print(f"✓ Presentation saved to: {output_path}")
        print(f"{'='*70}")
        
        # Show usage stats
        stats = self.llm_client.get_usage_stats()
        print(f"\nAI Usage Stats:")
        print(f"  Total tokens: {stats['total_tokens']:,}")
        print(f"  Total cost: ${stats['total_cost']:.4f}")
        print(f"  API calls: {stats['call_count']}")
        print()
        
        return str(output_path)
    
    def generate_from_file(
        self,
        input_file: str,
        output_path: Optional[str] = None,
        target_slides: Optional[int] = None
    ) -> str:
        """
        Generate presentation from text file using AI
        
        Args:
            input_file: Path to input text/markdown file
            output_path: Path to save output file
            target_slides: Optional target number of slides
            
        Returns:
            Path to generated presentation
        """
        input_path = Path(input_file)
        if not input_path.exists():
            raise FileNotFoundError(f"Input file not found: {input_file}")
        
        # Read content with encoding detection
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252', 'iso-8859-1']
        content = None
        
        for encoding in encodings:
            try:
                with open(input_path, 'r', encoding=encoding) as f:
                    content = f.read()
                break
            except UnicodeDecodeError:
                continue
        
        if content is None:
            raise ValueError(f"Could not decode file {input_file} with any supported encoding")
        
        # Generate presentation
        return self.generate_from_text(content, output_path, target_slides)


def main():
    """CLI entry point"""
    import sys
    
    if len(sys.argv) < 3:
        print("Usage: python smart_generator.py <template_path> <input_file> [output_file] [target_slides]")
        print("\nExample:")
        print("  python smart_generator.py templates/SavedTheme.pptx examples/sample.md")
        print("  python smart_generator.py templates/SavedTheme.pptx examples/sample.md output/result.pptx 12")
        print("\nNote: Requires OPENAI_API_KEY in .env file")
        sys.exit(1)
    
    template_path = sys.argv[1]
    input_file = sys.argv[2]
    output_file = sys.argv[3] if len(sys.argv) > 3 else None
    target_slides = int(sys.argv[4]) if len(sys.argv) > 4 else None
    
    try:
        generator = SmartGenerator(template_path)
        generator.generate_from_file(input_file, output_file, target_slides)
    
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
