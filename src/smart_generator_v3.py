"""
Smart Slide Generator V3 - Enhanced with schema-guided generation and holistic review
Two-stage planning + schema-guided content + batch review
"""

import sys
from pathlib import Path

# Add parent directory to path for imports when running as script
if __name__ == "__main__":
    sys.path.insert(0, str(Path(__file__).parent.parent))

import json
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
from typing import List, Dict, Optional, Any

from src.llm.client import LLMClient
from src.llm.enhanced_planner import EnhancedContentPlanner
from src.llm.schema_content_generator import SchemaGuidedGenerator
from src.llm.holistic_reviewer import HolisticReviewer
from src.core.slide_exporter import SlideExporter


class SmartGeneratorV3:
    """AI-powered slide generator with enhanced architecture"""
    
    def __init__(self, template_path: str, api_key: Optional[str] = None, schemas_path: str = "config/template_schemas.json"):
        """
        Initialize enhanced smart generator
        
        Args:
            template_path: Path to PowerPoint template file
            api_key: Optional OpenAI API key
            schemas_path: Path to template schemas JSON
        """
        self.template_path = Path(template_path)
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")
        
        # Initialize LLM components
        self.llm_client = LLMClient(api_key=api_key)
        self.planner = EnhancedContentPlanner(self.llm_client, schemas_path)
        self.content_generator = SchemaGuidedGenerator(self.llm_client)
        self.reviewer = HolisticReviewer(self.llm_client)
        self.exporter = SlideExporter()
        
        # Load template
        self.prs = Presentation(str(self.template_path))
        
        # Load schemas
        with open(schemas_path, 'r', encoding='utf-8') as f:
            schema_data = json.load(f)
            self.schemas = schema_data['schemas']
        
        # Get layout objects
        self.layout_dict = {}
        for layout in self.prs.slide_layouts:
            self.layout_dict[layout.name] = layout
        
        print(f"✓ Loaded template with {len(self.schemas)} layouts")
    
    def generate_from_text(
        self,
        text_content: str,
        output_path: Optional[str] = None,
        target_slides: Optional[int] = None,
        enable_review: bool = True
    ) -> str:
        """
        Generate presentation with enhanced process
        
        Args:
            text_content: Raw text/markdown content
            output_path: Path to save output file
            target_slides: Optional target number of slides
            enable_review: Whether to run Vision API review
            
        Returns:
            Path to generated presentation
        """
        print("\n" + "="*70)
        print("SMART SLIDE GENERATION V3 (Enhanced)")
        print("="*70 + "\n")
        
        # Clear existing template slides
        existing_count = len(self.prs.slides)
        if existing_count > 0:
            print(f"Removing {existing_count} existing slides from template...")
            for i in range(existing_count - 1, -1, -1):
                rId = self.prs.slides._sldIdLst[i].rId
                self.prs.part.drop_rel(rId)
                del self.prs.slides._sldIdLst[i]
        
        # STAGE 1: Create outline
        outline = self.planner.create_outline(
            content=text_content,
            target_slides=target_slides
        )
        
        # STAGE 2: Generate detailed content for each slide
        all_slide_content = self.content_generator.generate_all_slide_content(
            outline=outline,
            schemas=self.schemas,
            full_context=text_content
        )
        
        # STAGE 3: Build slides
        print(f"\nBuilding {len(all_slide_content)} slides...")
        for content in all_slide_content:
            self._create_slide_from_content(content)
        
        # Save presentation
        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = f"output/presentation_{timestamp}.pptx"
        
        output_file = Path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_file))
        
        print(f"\n✓ Presentation saved: {output_path}")
        
        # STAGE 4: Batch export and review (if enabled)
        if enable_review:
            review = self._run_holistic_review(output_file, text_content, outline)
            
            # Save review results
            review_path = output_file.with_suffix('.review.json')
            with open(review_path, 'w', encoding='utf-8') as f:
                json.dump(review, f, indent=2)
            print(f"\n✓ Review saved: {review_path}")
        
        # Print usage stats
        stats = self.llm_client.get_usage_stats()
        print(f"\n{'='*70}")
        print(f"✓ Presentation complete: {output_file}")
        print(f"{'='*70}\n")
        print(f"AI Usage Stats:")
        print(f"  Total tokens: {stats['total_tokens']:,}")
        print(f"  Total cost: ${stats['total_cost']:.4f}")
        print(f"  API calls: {stats['call_count']}")
        print()
        
        return str(output_file)
    
    def _create_slide_from_content(self, content: Dict[str, Any]):
        """Create a slide from generated content"""
        layout_name = content.get('layout_name')
        layout = self.layout_dict.get(layout_name)
        
        if not layout:
            print(f"  Warning: Layout '{layout_name}' not found, using default")
            layout = self.prs.slide_layouts[3]
        
        slide = self.prs.slides.add_slide(layout)
        
        # Populate fields based on what's available
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            # Title
            if shape.is_placeholder and 'TITLE' in str(shape.placeholder_format.type):
                if 'title' in content:
                    shape.text = content['title']
            
            # Subtitle
            elif shape.is_placeholder and shape.placeholder_format.idx == 1:
                if 'subtitle' in content:
                    shape.text = content['subtitle']
            
            # Content/bullets
            elif shape.is_placeholder and ('BODY' in str(shape.placeholder_format.type) or 'OBJECT' in str(shape.placeholder_format.type)):
                if 'content' in content and isinstance(content['content'], list):
                    text_frame = shape.text_frame
                    text_frame.clear()
                    
                    for bullet in content['content']:
                        p = text_frame.add_paragraph()
                        p.text = str(bullet)
                        p.level = 0
        
        # Add notes
        if 'notes' in content and content['notes']:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = content['notes']
    
    def _run_holistic_review(
        self,
        output_file: Path,
        original_content: str,
        outline: Dict
    ) -> Dict:
        """Export and review presentation"""
        print(f"\n{'='*70}")
        print("BATCH EXPORT & HOLISTIC REVIEW")
        print(f"{'='*70}")
        
        # Export all slides
        export_dir = output_file.parent / f"{output_file.stem}_slides"
        export_dir.mkdir(exist_ok=True)
        
        temp_exporter = SlideExporter(str(output_file))
        slide_images = temp_exporter.export_all_slides(str(export_dir))
        
        print(f"\n✓ Exported {len(slide_images)} slides to {export_dir}")
        
        # Review
        review = self.reviewer.review_presentation(
            slide_images=slide_images,
            original_content=original_content,
            outline=outline
        )
        
        return review
    
    def generate_from_file(
        self,
        input_file: str,
        output_path: Optional[str] = None,
        target_slides: Optional[int] = None,
        enable_review: bool = True
    ) -> str:
        """
        Generate presentation from file
        
        Args:
            input_file: Path to input text/markdown file
            output_path: Path to save output file
            target_slides: Optional target number of slides
            enable_review: Whether to run Vision review
            
        Returns:
            Path to generated presentation
        """
        input_path = Path(input_file)
        if not input_path.exists():
            raise FileNotFoundError(f"Input file not found: {input_file}")
        
        # Read content with encoding detection
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252', 'iso-8859-1']
        content = None
        
        for encoding in encodings:
            try:
                with open(input_path, 'r', encoding=encoding) as f:
                    content = f.read()
                break
            except UnicodeDecodeError:
                continue
        
        if content is None:
            raise ValueError(f"Could not decode file {input_file} with any supported encoding")
        
        # Generate presentation
        return self.generate_from_text(content, output_path, target_slides, enable_review)


def main():
    """CLI entry point"""
    import sys
    
    if len(sys.argv) < 3:
        print("Usage: python smart_generator_v3.py <template_path> <input_file> [output_file] [--no-review]")
        print("\nExample:")
        print("  python smart_generator_v3.py templates/SavedTheme.pptx examples/sample.md")
        print("  python smart_generator_v3.py templates/SavedTheme.pptx examples/sample.md output/result.pptx --no-review")
        print("\nNote: Requires OPENAI_API_KEY in .env file")
        sys.exit(1)
    
    template_path = sys.argv[1]
    input_file = sys.argv[2]
    output_file = sys.argv[3] if len(sys.argv) > 3 and not sys.argv[3].startswith('--') else None
    enable_review = '--no-review' not in sys.argv
    
    try:
        generator = SmartGeneratorV3(template_path)
        generator.generate_from_file(input_file, output_file, enable_review=enable_review)
    
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
