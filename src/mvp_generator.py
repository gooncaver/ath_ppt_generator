"""
MVP Slide Generator - Phase 1
Converts text content to basic PowerPoint slides without AI
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
import re
from datetime import datetime


class MVPGenerator:
    """Simple text-to-slides generator for Phase 1 MVP"""
    
    def __init__(self, template_path: str):
        """
        Initialize generator with template
        
        Args:
            template_path: Path to PowerPoint template file
        """
        self.template_path = Path(template_path)
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")
        
        self.prs = Presentation(str(self.template_path))
        self.default_layout_name = "10_Title and Content"
        
    def parse_markdown_content(self, content: str) -> list[dict]:
        """
        Parse markdown/text content into slide sections
        
        Args:
            content: Raw text content
            
        Returns:
            List of slide dicts with title and content
        """
        slides = []
        
        # Split by markdown headings (# Title)
        sections = re.split(r'\n(?=#+\s)', content)
        
        for section in sections:
            section = section.strip()
            if not section:
                continue
            
            # Extract title (first line if it's a heading)
            lines = section.split('\n')
            title = ""
            content_lines = []
            
            if lines[0].startswith('#'):
                # Remove markdown heading markers
                title = re.sub(r'^#+\s*', '', lines[0]).strip()
                content_lines = lines[1:]
            else:
                # Use first line as title, rest as content
                title = lines[0].strip()
                content_lines = lines[1:]
            
            # Process content
            content_text = '\n'.join(line.strip() for line in content_lines if line.strip())
            
            slides.append({
                'title': title,
                'content': content_text
            })
        
        return slides
    
    def get_layout_by_name(self, name: str):
        """
        Find layout by name in template
        
        Args:
            name: Layout name to search for
            
        Returns:
            Layout object or None
        """
        for layout in self.prs.slide_layouts:
            if layout.name == name:
                return layout
        return None
    
    def create_slide(self, title: str, content: str, layout_name: str = None):
        """
        Create a single slide with title and content
        
        Args:
            title: Slide title
            content: Slide content (will be split into bullets if multi-line)
            layout_name: Layout to use (defaults to Title and Content)
        """
        # Get layout
        if layout_name is None:
            layout_name = self.default_layout_name
        
        layout = self.get_layout_by_name(layout_name)
        if layout is None:
            # Fallback to first available layout with placeholders
            layout = self.prs.slide_layouts[1]
            print(f"Warning: Layout '{layout_name}' not found, using '{layout.name}'")
        
        # Create slide
        slide = self.prs.slides.add_slide(layout)
        
        # Populate title
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # Populate content
        # Find content placeholder (usually idx 1)
        content_placeholder = None
        for shape in slide.placeholders:
            # Look for body/content placeholder
            if shape.placeholder_format.idx == 1:
                content_placeholder = shape
                break
        
        if content_placeholder and hasattr(content_placeholder, 'text_frame'):
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # Clear default text
            
            # Split content into lines/bullets
            lines = [line.strip() for line in content.split('\n') if line.strip()]
            
            if lines:
                # First line
                p = text_frame.paragraphs[0]
                p.text = lines[0]
                p.level = 0
                
                # Additional lines as separate paragraphs
                for line in lines[1:]:
                    p = text_frame.add_paragraph()
                    # Handle markdown bullets
                    if line.startswith('- ') or line.startswith('* '):
                        p.text = line[2:]
                    else:
                        p.text = line
                    p.level = 0
    
    def generate_from_text(self, text_content: str, output_path: str = None):
        """
        Generate presentation from text content
        
        Args:
            text_content: Raw text/markdown content
            output_path: Path to save output file (auto-generated if None)
        """
        # Parse content into slides
        slide_data = self.parse_markdown_content(text_content)
        
        if not slide_data:
            print("Warning: No content found to generate slides")
            return
        
        # Clear any existing slides from template
        # Remove slides in reverse order to avoid index issues
        existing_count = len(self.prs.slides)
        if existing_count > 0:
            print(f"Removing {existing_count} existing slides from template...")
            for i in range(existing_count - 1, -1, -1):
                rId = self.prs.slides._sldIdLst[i].rId
                self.prs.part.drop_rel(rId)
                del self.prs.slides._sldIdLst[i]
        
        print(f"Generating {len(slide_data)} slides...")
        
        # Create slides
        for idx, slide_info in enumerate(slide_data, 1):
            print(f"  [{idx}/{len(slide_data)}] {slide_info['title']}")
            self.create_slide(slide_info['title'], slide_info['content'])
        
        # Save presentation
        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = f"output/presentation_{timestamp}.pptx"
        
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        self.prs.save(str(output_path))
        print(f"\n✓ Presentation saved to: {output_path}")
        
        return str(output_path)
    
    def generate_from_file(self, input_file: str, output_path: str = None):
        """
        Generate presentation from text file
        
        Args:
            input_file: Path to input text/markdown file
            output_path: Path to save output file
        """
        input_path = Path(input_file)
        if not input_path.exists():
            raise FileNotFoundError(f"Input file not found: {input_file}")
        
        # Read content
        with open(input_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Generate presentation
        return self.generate_from_text(content, output_path)


def main():
    """CLI entry point"""
    import sys
    
    if len(sys.argv) < 3:
        print("Usage: python mvp_generator.py <template_path> <input_file> [output_file]")
        print("\nExample:")
        print("  python mvp_generator.py templates/base_template.pptx examples/sample.md")
        print("  python mvp_generator.py templates/base_template.pptx examples/sample.txt output/result.pptx")
        sys.exit(1)
    
    template_path = sys.argv[1]
    input_file = sys.argv[2]
    output_file = sys.argv[3] if len(sys.argv) > 3 else None
    
    try:
        generator = MVPGenerator(template_path)
        generator.generate_from_file(input_file, output_file)
    
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
