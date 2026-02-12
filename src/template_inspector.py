"""
Template Inspector - Analyzes PowerPoint template to discover available layouts
Phase 1 MVP Tool
"""

from pptx import Presentation
from pathlib import Path
import sys


class TemplateInspector:
    """Inspects PowerPoint templates to extract layout information"""
    
    def __init__(self, template_path: str):
        """
        Initialize inspector with template path
        
        Args:
            template_path: Path to .pptx or .potx template file
        """
        self.template_path = Path(template_path)
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")
        
        self.prs = Presentation(str(self.template_path))
        
    def get_slide_layouts(self) -> list[dict]:
        """
        Extract all available slide layouts from template
        
        Returns:
            List of dicts with layout info: name, index, placeholder count
        """
        layouts = []
        
        for idx, layout in enumerate(self.prs.slide_layouts):
            layout_info = {
                'index': idx,
                'name': layout.name,
                'placeholder_count': len(layout.placeholders),
                'placeholders': []
            }
            
            # Extract placeholder information
            for ph_idx, placeholder in enumerate(layout.placeholders):
                ph_info = {
                    'idx': placeholder.placeholder_format.idx,
                    'type': placeholder.placeholder_format.type,
                    'name': placeholder.name
                }
                layout_info['placeholders'].append(ph_info)
            
            layouts.append(layout_info)
        
        return layouts
    
    def print_layouts(self):
        """Print formatted layout information to console"""
        layouts = self.get_slide_layouts()
        
        print("=" * 80)
        print(f"TEMPLATE: {self.template_path.name}")
        print("=" * 80)
        print(f"\nTotal Layouts: {len(layouts)}\n")
        
        for layout in layouts:
            print(f"[{layout['index']}] {layout['name']}")
            print(f"    Placeholders: {layout['placeholder_count']}")
            
            if layout['placeholders']:
                for ph in layout['placeholders']:
                    print(f"      - {ph['name']} (idx={ph['idx']}, type={ph['type']})")
            print()
    
    def save_layout_info(self, output_path: str):
        """
        Save layout information to text file
        
        Args:
            output_path: Path to output file
        """
        layouts = self.get_slide_layouts()
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write("=" * 80 + "\n")
            f.write(f"TEMPLATE: {self.template_path.name}\n")
            f.write("=" * 80 + "\n")
            f.write(f"\nTotal Layouts: {len(layouts)}\n\n")
            
            for layout in layouts:
                f.write(f"[{layout['index']}] {layout['name']}\n")
                f.write(f"    Placeholders: {layout['placeholder_count']}\n")
                
                if layout['placeholders']:
                    for ph in layout['placeholders']:
                        f.write(f"      - {ph['name']} (idx={ph['idx']}, type={ph['type']})\n")
                f.write("\n")
        
        print(f"✓ Layout information saved to: {output_path}")


def main():
    """CLI entry point for template inspector"""
    if len(sys.argv) < 2:
        print("Usage: python template_inspector.py <template_path> [output_file]")
        print("\nExample:")
        print("  python template_inspector.py templates/base_template.pptx")
        print("  python template_inspector.py templates/base_template.pptx docs/layouts.txt")
        sys.exit(1)
    
    template_path = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None
    
    try:
        inspector = TemplateInspector(template_path)
        inspector.print_layouts()
        
        if output_file:
            inspector.save_layout_info(output_file)
    
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
